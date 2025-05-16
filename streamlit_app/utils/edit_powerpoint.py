import pandas as pd
import numpy as np
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
import os
import io
import logging
from typing import Dict, List, Tuple, Optional, Union, Any, IO

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger('powerpoint_editor')

# Define constants for better maintainability
class TarifaConstants:
    # ENEL CE Tariffs
    TE_ENEL_CE = 0.27291
    TUSD_ENEL_CE = 0.44929
    
    # Tax rates
    ICMS = 0.2
    PISCOF = 0.05
    
    @classmethod
    def calculate_derived_rates(cls):
        """Calculate derived tariff rates based on base constants"""
        # Base tariff without taxes
        tarifa_de_aplicacao = cls.TE_ENEL_CE + cls.TUSD_ENEL_CE
        
        # Tariff with taxes
        tarifa_fornecida = round(tarifa_de_aplicacao / ((1 - cls.ICMS) * (1 - cls.PISCOF)), 4)
        tarifa_injetada_compartilhada = cls.TE_ENEL_CE / ((1 - cls.ICMS) * (1 - cls.PISCOF)) + cls.TUSD_ENEL_CE / (1 - cls.PISCOF)
        tarifa_injetada_real = round(tarifa_injetada_compartilhada, 4)
        impostos_n_compoensaveis = tarifa_fornecida - tarifa_injetada_real
        
        return {
            'tarifa de aplicação': tarifa_de_aplicacao,
            'tarifa fornecida': tarifa_fornecida,
            'tarifa injetada compartilhada': tarifa_injetada_compartilhada,
            'tarifa injetada real': tarifa_injetada_real,
            'impostos não compensáveis': impostos_n_compoensaveis
            }
    
    def get_tarifas(self):
        tarifas = self.calculate_derived_rates()

        fornecida = tarifas['tarifa fornecida']
        injetada = tarifas['tarifa injetada compartilhada']
        impostos = tarifas['impostos não compensáveis']
        injetada_real = tarifas['tarifa injetada real']
        return fornecida, injetada, impostos, injetada_real



class TextFormatter:
    """Class for text formatting operations in PowerPoint"""
    
    @staticmethod
    def format_text(run, text: str, font_name: str, font_size: int, align_center: bool = False):
        """Format text run with specified parameters"""
        run.text = text
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        if align_center:
            run.alignment = PP_ALIGN.CENTER

    @staticmethod
    def format_money_br(value: float) -> str:
        """Format monetary values in Brazilian format (R$ 1.234,56)"""
        return f"{value:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

    @staticmethod
    def update_text_shape(shape, text1: str, font1: str, size1: int, 
                         text2: Optional[str] = None, font2: Optional[str] = None, 
                         size2: Optional[int] = None, center_align: bool = True):
        """Update text in a shape with proper formatting"""
        # Set defaults for second text if not provided
        if font2 is None:
            font2 = font1
        if size2 is None:
            size2 = size1

        # Get the TextFrame from the shape
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        # Update the primary text with existing formatting
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            if center_align:
                p.alignment = PP_ALIGN.CENTER
                
            if p.runs:
                run = p.runs[0]
                TextFormatter.format_text(run, text1, font1, size1, center_align)
            else:
                p.text = text1
                p.font.name = font1
                p.font.size = Pt(size1)
                if center_align:
                    p.alignment = PP_ALIGN.CENTER

        # Update secondary text if provided
        if text2:
            if len(text_frame.paragraphs) > 1:
                p = text_frame.paragraphs[1]
                if center_align:
                    p.alignment = PP_ALIGN.CENTER
                    
                if p.runs:
                    run = p.runs[0]
                    TextFormatter.format_text(run, text2, font2, size2, center_align)
                else:
                    p.text = text2
                    p.font.name = font2
                    p.font.size = Pt(size2)
                    if center_align:
                        p.alignment = PP_ALIGN.CENTER
            else:
                p = text_frame.add_paragraph()
                p.text = text2
                p.font.name = font2
                p.font.size = Pt(size2)
                if center_align:
                    p.alignment = PP_ALIGN.CENTER


class ChartUpdater:
    """Class for chart updating operations in PowerPoint"""
    
    @staticmethod
    def update_chart(chart, categories: List[str], *series: Tuple[str, List[float]]):
        """Update chart with categories and data series"""
        try:
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for name, values in series:
                chart_data.add_series(name, values)
            chart.replace_data(chart_data)
            
            # Update formatting
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(10)  # Default size
            category_axis.tick_labels.font.name = "Calibri"  # Default font
            
            return True
        except Exception as e:
            logger.error(f"Error updating chart: {str(e)}")
            return False


class FinancialCalculator:
    """Handles financial calculations for the proposal"""
    
    def __init__(self, tariff_data: TarifaConstants):
        """Initialize with tariff data"""
        self.tariff_data = tariff_data
    
    def calculate_bill_before(self, pic: float, disponibility_cost: str, energy_consumption: float, n_ucs: int) -> Dict[str, float]:
        """Calculate the bill before the proposal"""
        tarifas_fornecida, tarifa_injetada, tarifas_impostos, tarifa_injetada_real = self.tariff_data.get_tarifas()

        energy_cost = (energy_consumption - (n_ucs * 100)) * tarifa_injetada_real

        if disponibility_cost == 'Trifásico':
            min_cost = 100 * n_ucs * tarifas_fornecida
        elif disponibility_cost == 'Monofásico':
            min_cost = 30 * n_ucs * tarifas_fornecida
        elif disponibility_cost == 'Bifásico':
            min_cost = 50 * n_ucs * tarifas_fornecida
        else:
            raise ValueError("Custo de disponibilidade inválido.")
        
        public_ilumination_cost = pic
        
        taxes = (energy_consumption - n_ucs * 100) * tarifas_impostos
        
        total_bill = energy_cost + min_cost + public_ilumination_cost + taxes
        return {
            'valor em energia': energy_cost,
            'custo mínimo': min_cost,
            'iluminação pública': public_ilumination_cost,
            'impostos': taxes,
            'total': total_bill
        }

    def calculate_bill_after(self, pic: float, disponibility_cost: str, energy_consumption: float, n_ucs: int, bill_discount: float) -> Dict[str, float]:
        """Calculate the bill after the proposal"""
        enel_bill = self.calculate_bill_before(pic, disponibility_cost, energy_consumption, n_ucs)
        
        taxes = enel_bill['impostos']
        energy_cost = enel_bill['valor em energia']
        min_cost = enel_bill['custo mínimo']
        public_ilumination_cost = enel_bill['iluminação pública']

        energy_cost_discounted = energy_cost * (1 - (bill_discount/100))

        total_bill_discounted = taxes + public_ilumination_cost + min_cost + energy_cost_discounted       
        
        return {
            'valor em energia': energy_cost,
            'custo mínimo': min_cost,
            'iluminação pública': public_ilumination_cost,
            'impostos': taxes,
            'total': total_bill_discounted
        }


class PowerPointUpdater:
    """Main class for PowerPoint presentation updates"""
    
    def __init__(self, template_buffer: IO):
        """Initialize with template file buffer"""
        self.presentation = Presentation(template_buffer)
        self.text_formatter = TextFormatter()
    
    def update_presentation(self, client_info: Dict[str, Any], conta_antes: Dict[str, float], conta_depois: Dict[str, float]) -> None:
        """Update the entire presentation with client info and financial data"""
        
        # Extract client info
        client = client_info['cliente']
        discount = client_info['desconto']
        
        # Update slides
        for slide in self.presentation.slides:
            self._update_slide(
                slide, 
                client, 
                discount, 
                conta_antes,
                conta_depois
            )
    
    def _update_slide(self, 
                     slide, 
                     client: str, 
                     discount: float, 
                     conta_antes: Dict[str, float],
                     conta_depois: Dict[str, float]) -> None:
        """Update a single slide with all relevant information"""
        
        for shape in slide.shapes:
            # Only process text boxes
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    self._update_text_shape(
                        shape, 
                        text, 
                        client, 
                        discount,
                        conta_antes,
                        conta_depois
                    )
    
    def _update_text_shape(self, 
                          shape, 
                          text: str, 
                          client: str, 
                          discount: float,
                          conta_antes: Dict[str, float],
                          conta_depois: Dict[str, float]) -> None:
        """Update text content in a shape based on its current text"""
        
        try:
            if text.startswith("CLIENTE: PPPPPP"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"CLIENTE: {client}", 
                    font1="Arial", 
                    size1=24, 
                    center_align=False
                )

            if text.startswith("DATA: DDDDDDD"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"DATA: {datetime.datetime.now().strftime('%d/%m/%Y')}", 
                    font1="Arial", 
                    size1=24, 
                    center_align=False
                )
                
            elif text.startswith("XX%"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"{discount}%", 
                    font1="Inter Bold", 
                    size1=20
                )

            elif text.startswith("YY%"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"{discount}%", 
                    font1="Inter Bold", 
                    size1=21
                )
                            
            # Current consumption cost field
            elif text.startswith("R$ AAAA"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['total'])}", 
                    font1="Inter Bold", 
                    size1=26
                )

            elif text.startswith("R$ AAAB"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['total'])}", 
                    font1="Inter Bold", 
                    size1=21
                )
                
            elif text.startswith("R$ BBBB"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['custo mínimo'])}", 
                    font1="Inter Bold", 
                    size1=21
                )

            elif text.startswith(" R$ aBBa"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['custo mínimo'])}", 
                    font1="Inter", 
                    size1=21
                )
                
            elif text.startswith("R$ CCCC"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_depois['valor em energia'])}", 
                    font1="Inter Bold", 
                    size1=21
                )
                
            elif text.startswith("R$ DDDD"):
                savings = conta_antes['total'] - conta_depois['total']
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(savings)}", 
                    font1="Inter Bold", 
                    size1=21
                )

            elif text.startswith("R$ DDDB"):
                savings = conta_antes['total'] - conta_depois['total']
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(savings)}", 
                    font1="Inter Bold", 
                    size1=21
                ) 

            elif text.startswith("R$ EEEE"):
                savings = conta_antes['total'] - conta_depois['total']
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(savings)}", 
                    font1="Inter Bold", 
                    size1= 37
                )

            elif text.startswith("R$ EEEB"):
                savings = conta_antes['total'] - conta_depois['total']
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(savings)}", 
                    font1="Inter Bold", 
                    size1= 21
                ) 

            elif text.startswith("R$ FFFF"):
                savings = conta_antes['total'] - conta_depois['total']
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(12 * savings)}", 
                    font1="Inter Bold", 
                    size1=37
                )

            elif text.startswith("R$ GGGG"):
                savings = conta_antes['total'] - conta_depois['total']
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(12 * 5 * savings)}", 
                    font1="Inter Bold", 
                    size1=37
                ) 
                
            elif text.startswith("R$ HHHH"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['valor em energia'])}", 
                    font1="Inter", 
                    size1=21
                ) 
            
            elif text.startswith("R$ IIII"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['iluminação pública'])}", 
                    font1="Inter", 
                    size1=21
                ) 
            
            elif text.startswith("R$ CDCD"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['custo mínimo'])}", 
                    font1="Inter", 
                    size1=21
                )

            elif text.startswith("R$ CICI"):
                TextFormatter.update_text_shape(
                    shape, 
                    text1=f"R$ {TextFormatter.format_money_br(conta_antes['impostos'])}", 
                    font1="Inter", 
                    size1=21
                ) 
             
            
        
        except Exception as e:
            logger.error(f"Error updating text shape with text '{text[:20]}...': {str(e)}")
    
    def save_presentation(self) -> io.BytesIO:
        """Save the presentation to a BytesIO buffer and return it"""
        output_buffer = io.BytesIO()
        self.presentation.save(output_buffer)
        output_buffer.seek(0)  # Reset buffer position to beginning
        return output_buffer


def powerpoint_edit(infos: Dict[str, Any], buffer: IO) -> io.BytesIO:
    """Main function to edit PowerPoint with client information
    
    Args:
        infos: Dictionary with client information
        buffer: Input file buffer containing the PowerPoint template
        
    Returns:
        BytesIO buffer containing the edited PowerPoint file
    """
    try:
        # Calculate tariff constants
        tariff_data = TarifaConstants()
        
        # Initialize financial calculator
        calculator = FinancialCalculator(tariff_data)
        
        # Calculate financial values
        conta_antes = calculator.calculate_bill_before(
            infos['cip'], 
            infos['custo_disponibilidade'], 
            infos['consumo'], 
            infos['n_ucs']
        )

        conta_depois = calculator.calculate_bill_after(
            infos['cip'], 
            infos['custo_disponibilidade'], 
            infos['consumo'], 
            infos['n_ucs'], 
            infos['desconto']
        )
        
        # Initialize PowerPoint updater
        ppt_updater = PowerPointUpdater(buffer)
        
        # Update presentation with client info and financial data
        ppt_updater.update_presentation(infos, conta_antes, conta_depois)
        
        # Save and return the updated presentation
        return ppt_updater.save_presentation()
        
    except Exception as e:
        logger.error(f"Error in powerpoint_edit: {str(e)}")
        raise


# Optional PDF conversion function (commented out in original code)
'''
def convert_ppt_to_pdf(input_file, output_file):
    """Convert PowerPoint file to PDF using COM automation
    
    This requires Windows with PowerPoint installed
    
    Args:
        input_file: Path to input PowerPoint file
        output_file: Path for output PDF file
    """
    try:
        import comtypes.client
        
        # Initialize COM
        comtypes.CoInitialize()
        
        # Convert to absolute paths
        input_file = os.path.abspath(input_file)
        output_file = os.path.abspath(output_file)
        
        # Open PowerPoint
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.DisplayAlerts = 0  # Disable alerts
        
        # Open presentation and save as PDF
        presentation = powerpoint.Presentations.Open(input_file, WithWindow=False)
        presentation.SaveAs(output_file, 32)  # 32 is the format type for PDF
        
        # Clean up
        presentation.Close()
        powerpoint.Quit()
        comtypes.CoUninitialize()
        
        return True
    except Exception as e:
        logger.error(f"Error converting PowerPoint to PDF: {str(e)}")
        return False
'''